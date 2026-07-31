from __future__ import annotations

import sys
import unittest
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"

if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pptx_quicktool.page_plan import generate_page_plan
from pptx_quicktool.pptx_exporter import DEFAULT_TEMPLATE_PATH, export_page_plan_to_pptx
from pptx_quicktool.pptx_verifier import linked_slide_indices
from pptx_quicktool.training_document import SectionInput, TrainingDocumentInput


def slide_texts(path: Path) -> list[list[str]]:
    presentation = Presentation(path)
    slides: list[list[str]] = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().replace("\v", "\n"))
        slides.append(texts)
    return slides


def shape_with_text(slide, text: str):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip().replace("\v", "\n") == text:
            return shape
    raise AssertionError(f"Could not find shape with text: {text}")


def slide_combined_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().replace("\v", "\n"))
    return "\n".join(texts)


def paragraph_with_text(slide, text: str):
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text == text:
                return paragraph
    raise AssertionError(f"Could not find paragraph with text: {text}")


def agenda_shapes(slide):
    return [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text.strip() and shape.text.strip() != "目錄"
    ]


def body_placeholders(slide):
    return [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    ]


def placeholder_shape_with_text(slide, text: str):
    shape = shape_with_text(slide, text)
    if not shape.is_placeholder:
        raise AssertionError(f"Shape with text is not a placeholder: {text}")
    return shape


def title_placeholder(layout):
    for placeholder in layout.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return placeholder
    raise AssertionError("Layout has no title placeholder")


class PptxExporterTests(unittest.TestCase):
    def test_exports_valid_pptx_with_expected_slides_and_text(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[
                SectionInput(
                    title="Setup",
                    content_page_titles=["Install Tools", "Connect Controller"],
                ),
                SectionInput(title="Operation", content_page_titles=["Open Project"]),
            ],
        )
        plan = generate_page_plan(document)

        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document.pptx"
        if output_path.exists():
            output_path.unlink()

        export_page_plan_to_pptx(plan, output_path)

        self.assertTrue(output_path.exists())
        presentation = Presentation(output_path)
        self.assertEqual(len(presentation.slides), 7)
        self.assertIn("TwinCAT Training", slide_combined_text(presentation.slides[0]))
        self.assertIn("目錄", slide_combined_text(presentation.slides[1]))
        self.assertIn("Setup", slide_combined_text(presentation.slides[1]))
        self.assertIn("Install Tools", slide_combined_text(presentation.slides[1]))
        self.assertIn("Connect Controller", slide_combined_text(presentation.slides[1]))
        self.assertIn("Operation", slide_combined_text(presentation.slides[1]))
        self.assertIn("Setup", slide_combined_text(presentation.slides[2]))
        self.assertIn("Install Tools", slide_combined_text(presentation.slides[3]))
        self.assertNotIn("回主目錄", slide_combined_text(presentation.slides[3]))

    def test_exports_table_of_contents_and_return_navigation_links(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[
                SectionInput(
                    title="Setup",
                    content_page_titles=["Install Tools", "Connect Controller"],
                ),
                SectionInput(title="Operation", content_page_titles=["Open Project"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-links.pptx"
        if output_path.exists():
            output_path.unlink()

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertEqual(linked_slide_indices(presentation, presentation.slides[1]), [2, 5])
        self.assertEqual(linked_slide_indices(presentation, presentation.slides[2]), [1])
        self.assertIn("目錄", slide_combined_text(presentation.slides[2]))
        self.assertIn("Operation", slide_combined_text(presentation.slides[2]))
        self.assertEqual(linked_slide_indices(presentation, presentation.slides[5]), [1])
        for slide_index in [3, 4, 6]:
            self.assertEqual(linked_slide_indices(presentation, presentation.slides[slide_index]), [])

    def test_table_of_contents_links_follow_reordered_sections(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[
                SectionInput(title="Operation", content_page_titles=["Open Project"]),
                SectionInput(
                    title="Setup",
                    content_page_titles=["Install Tools", "Connect Controller"],
                ),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-reordered-links.pptx"
        if output_path.exists():
            output_path.unlink()

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertEqual(linked_slide_indices(presentation, presentation.slides[1]), [2, 4])

    def test_agenda_text_leaves_bullets_and_numbering_to_the_template(self) -> None:
        document = TrainingDocumentInput(
            title="測試文件",
            sections=[
                SectionInput(title="測試1", content_page_titles=["1", "2"]),
                SectionInput(title="測試2", content_page_titles=["1", "2"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-template-agenda-text.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        table_of_contents_text = slide_combined_text(presentation.slides[1])
        section_text = slide_combined_text(presentation.slides[2])

        self.assertIn("測試1", table_of_contents_text)
        self.assertIn("測試2", table_of_contents_text)
        self.assertIn("1", table_of_contents_text)
        self.assertNotIn("1. 測試1", table_of_contents_text)
        self.assertNotIn("2. 測試2", table_of_contents_text)
        self.assertNotIn("- 1", table_of_contents_text)
        self.assertNotIn("- 2", table_of_contents_text)
        self.assertNotIn("1. 測試1", section_text)
        self.assertNotIn("- 1", section_text)

    def test_agenda_paragraphs_use_template_bullets_instead_of_auto_numbering(self) -> None:
        document = TrainingDocumentInput(
            title="測試文件",
            sections=[
                SectionInput(title="測試1", content_page_titles=["1", "2"]),
                SectionInput(title="測試2", content_page_titles=["1", "2"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-template-bullets.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        section_paragraph = paragraph_with_text(presentation.slides[1], "測試1")
        content_paragraph = paragraph_with_text(presentation.slides[1], "1")

        self.assertEqual(section_paragraph.level, 0)
        self.assertEqual(content_paragraph.level, 1)
        self.assertEqual(section_paragraph._p.pPr.xpath("./a:buAutoNum"), [])
        self.assertEqual(content_paragraph._p.pPr.xpath("./a:buAutoNum"), [])
        self.assertTrue(section_paragraph._p.pPr.xpath("./a:buChar"))

    def test_exports_reference_size_and_consistent_title_hierarchy(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[
                SectionInput(title="Setup", content_page_titles=["Install Tools"]),
                SectionInput(title="Operation", content_page_titles=["Open Project"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-format.pptx"
        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertEqual(presentation.slide_width, 12192000)
        self.assertEqual(presentation.slide_height, 6858000)

        title_shapes = [
            placeholder_shape_with_text(presentation.slides[0], "TwinCAT Training"),
            placeholder_shape_with_text(presentation.slides[1], "\u76ee\u9304"),
            placeholder_shape_with_text(presentation.slides[2], "Setup"),
            placeholder_shape_with_text(presentation.slides[3], "Install Tools"),
        ]
        template = Presentation(DEFAULT_TEMPLATE_PATH)
        template_title_placeholder = title_placeholder(template.slide_layouts[0])
        self.assertEqual({shape.left for shape in title_shapes}, {template_title_placeholder.left})
        self.assertEqual({shape.top for shape in title_shapes}, {template_title_placeholder.top})
        self.assertEqual({shape.text_frame.paragraphs[0].runs[0].font.name for shape in title_shapes}, {None})
        self.assertEqual({shape.text_frame.paragraphs[0].runs[0].font.size for shape in title_shapes}, {None})

    def test_exports_with_default_training_document_template_layouts(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[SectionInput(title="Setup", content_page_titles=["Install Tools"])],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-default-template.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertTrue(DEFAULT_TEMPLATE_PATH.exists())
        self.assertEqual(
            [slide.slide_layout.name for slide in presentation.slides],
            ["Title with picture", "Contents", "Contents", "Text"],
        )

    def test_exports_with_specified_template_path_without_modifying_source(self) -> None:
        document = TrainingDocumentInput(
            title="TwinCAT Training",
            sections=[SectionInput(title="Setup", content_page_titles=["Install Tools"])],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-specified-template.pptx"
        template_mtime = DEFAULT_TEMPLATE_PATH.stat().st_mtime_ns

        export_page_plan_to_pptx(plan, output_path, template_path=DEFAULT_TEMPLATE_PATH)

        presentation = Presentation(output_path)
        self.assertEqual(len(presentation.slides), 4)
        self.assertEqual(presentation.slides[0].slide_layout.name, "Title with picture")
        self.assertEqual(DEFAULT_TEMPLATE_PATH.stat().st_mtime_ns, template_mtime)

    def test_agenda_wraps_to_at_most_three_columns_before_exceeding_the_body_area(self) -> None:
        document = TrainingDocumentInput(
            title="測試文件",
            sections=[
                SectionInput(title=f"章節{i}", content_page_titles=[str(page) for page in range(1, 6)])
                for i in range(1, 8)
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-multi-column-agenda.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        table_of_contents_shapes = agenda_shapes(presentation.slides[1])

        self.assertEqual(len(table_of_contents_shapes), 3)
        self.assertEqual(sorted({shape.left for shape in table_of_contents_shapes}), [shape.left for shape in table_of_contents_shapes])
        for shape in table_of_contents_shapes:
            populated_paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
            self.assertLessEqual(len(populated_paragraphs), 15)
            self.assertLessEqual(shape.top + shape.height, presentation.slide_height)

    def test_exports_section_without_content_pages(self) -> None:
        document = TrainingDocumentInput(
            title="測試文件",
            sections=[
                SectionInput(title="前言", content_page_titles=[]),
                SectionInput(title="章節1", content_page_titles=["1"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-section-without-content.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertEqual(len(presentation.slides), 6)
        self.assertIn("前言", slide_combined_text(presentation.slides[1]))
        self.assertIn("前言", slide_combined_text(presentation.slides[2]))
        self.assertIn("前言", slide_combined_text(presentation.slides[3]))
        self.assertEqual(presentation.slides[3].slide_layout.name, "Text")
        self.assertEqual(linked_slide_indices(presentation, presentation.slides[1]), [2, 4])

    def test_agenda_slides_remove_template_body_placeholders_after_populating_columns(self) -> None:
        document = TrainingDocumentInput(
            title="測試文件",
            sections=[
                SectionInput(title="前言", content_page_titles=[]),
                SectionInput(title="章節1", content_page_titles=["1"]),
            ],
        )
        plan = generate_page_plan(document)
        output_dir = ROOT / ".tmp" / "test-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "training-document-no-agenda-placeholder.pptx"

        export_page_plan_to_pptx(plan, output_path)

        presentation = Presentation(output_path)
        self.assertEqual(body_placeholders(presentation.slides[1]), [])
        self.assertEqual(body_placeholders(presentation.slides[2]), [])


if __name__ == "__main__":
    unittest.main()
