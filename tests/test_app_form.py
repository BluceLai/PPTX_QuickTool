from __future__ import annotations

import sys
import unittest
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"

if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from pptx_quicktool.app import create_main_window
from pptx import Presentation


def slide_texts(path: Path) -> list[list[str]]:
    presentation = Presentation(path)
    slides: list[list[str]] = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().replace("\v", "\n"))
        slides.append(texts)
    return slides


def slide_combined_texts(path: Path) -> list[str]:
    return ["\n".join(texts) for texts in slide_texts(path)]


class AppFormTests(unittest.TestCase):
    def test_form_edits_document_structure_and_preview(self) -> None:
        root = create_main_window()
        try:
            form = root.form

            form.set_document_title("TwinCAT Training")
            form.add_section("Setup")
            form.add_content_page("Install Tools")
            form.add_content_page("Connect Controller")
            form.add_section("Operation")
            form.add_content_page("Open Project")

            document = form.current_document()
            self.assertEqual(document.title, "TwinCAT Training")
            self.assertEqual([section.title for section in document.sections], ["Setup", "Operation"])
            self.assertEqual(document.sections[0].content_page_titles, ["Install Tools", "Connect Controller"])
            self.assertEqual(document.sections[1].content_page_titles, ["Open Project"])
            self.assertEqual(form.validation_messages(), [])
            self.assertEqual(
                form.preview_text(),
                (
                    "1. \u5c01\u9762 - TwinCAT Training\n"
                    "2. \u76ee\u9304\n"
                    "3. \u7ae0\u7bc0 - Setup\n"
                    "4. \u5167\u6587 - Install Tools\n"
                    "5. \u5167\u6587 - Connect Controller\n"
                    "6. \u7ae0\u7bc0 - Operation\n"
                    "7. \u5167\u6587 - Open Project"
                ),
            )
        finally:
            root.update()
            root.destroy()

    def test_form_renames_removes_and_reorders_sections_and_pages(self) -> None:
        root = create_main_window()
        try:
            form = root.form

            form.set_document_title("TwinCAT Training")
            form.add_section("Setup")
            form.add_content_page("Install Tools")
            form.add_content_page("Connect Controller")
            form.add_section("Operation")
            form.add_content_page("Open Project")

            form.select_section(1)
            form.rename_selected_section("Operate")
            form.move_selected_section_up()
            form.select_section(0)
            form.rename_content_page(0, "Open Existing Project")
            form.add_content_page("Run Check")
            form.move_content_page_down(0)
            form.remove_content_page(0)

            document = form.current_document()
            self.assertEqual([section.title for section in document.sections], ["Operate", "Setup"])
            self.assertEqual(document.sections[0].content_page_titles, ["Open Existing Project"])
            self.assertEqual(document.sections[1].content_page_titles, ["Install Tools", "Connect Controller"])
        finally:
            root.update()
            root.destroy()

    def test_form_shows_validation_messages_from_input_model(self) -> None:
        root = create_main_window()
        try:
            form = root.form

            form.set_document_title("")
            form.add_section("")

            self.assertEqual(
                form.validation_messages(),
                [
                    "Document title is required.",
                    "Section 1 title is required.",
                    "Section 1 must include at least one content page.",
                ],
            )
            self.assertEqual(
                form.validation_var.get(),
                (
                    "\u8acb\u8f38\u5165 PPT \u6a19\u984c\u3002 "
                    "\u7b2c 1 \u500b\u7ae0\u7bc0\u9700\u8981\u6a19\u984c\u3002 "
                    "\u7b2c 1 \u500b\u7ae0\u7bc0\u81f3\u5c11\u9700\u8981\u4e00\u500b\u5167\u6587\u9801\u3002"
                ),
            )
            self.assertEqual(form.preview_text(), "")
        finally:
            root.update()
            root.destroy()

    def test_form_generates_pptx_from_current_data(self) -> None:
        root = create_main_window()
        try:
            form = root.form
            form.set_document_title("TwinCAT Training")
            form.add_section("Setup")
            form.add_content_page("Install Tools")
            form.add_section("Operation")
            form.add_content_page("Open Project")

            output_dir = ROOT / ".tmp" / "test-output"
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / "ui-generated.pptx"
            if output_path.exists():
                output_path.unlink()

            result = form.generate_to_path(output_path)

            self.assertEqual(result, output_path)
            self.assertTrue(output_path.exists())
            texts = slide_combined_texts(output_path)
            self.assertIn("TwinCAT Training", texts[0])
            self.assertIn("目錄", texts[1])
            self.assertIn("Setup", texts[1])
            self.assertIn("Install Tools", texts[1])
            self.assertIn("Operation", texts[1])
            self.assertIn("目錄", texts[2])
            self.assertIn("Open Project", texts[4])
            self.assertNotIn("回主目錄", texts[3])
            self.assertNotIn("回主目錄", texts[5])
            self.assertEqual(form.generation_status_var.get(), f"\u5df2\u7522\u751f\u4e26\u9a57\u8b49\uff1a{output_path}")
        finally:
            root.update()
            root.destroy()

    def test_form_blocks_generation_when_validation_fails(self) -> None:
        root = create_main_window()
        try:
            form = root.form
            output_dir = ROOT / ".tmp" / "test-output"
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / "invalid-ui-generated.pptx"
            if output_path.exists():
                output_path.unlink()

            result = form.generate_to_path(output_path)

            self.assertIsNone(result)
            self.assertFalse(output_path.exists())
            self.assertIn("\u8acb\u8f38\u5165 PPT \u6a19\u984c\u3002", form.generation_status_var.get())
        finally:
            root.update()
            root.destroy()

    def test_form_reports_generation_errors(self) -> None:
        root = create_main_window()
        try:
            form = root.form
            form.set_document_title("TwinCAT Training")
            form.add_section("Setup")
            form.add_content_page("Install Tools")
            output_dir = ROOT / ".tmp" / "test-output"
            output_dir.mkdir(parents=True, exist_ok=True)

            result = form.generate_to_path(output_dir)

            self.assertIsNone(result)
            self.assertTrue(form.generation_status_var.get().startswith("\u7522\u751f\u5931\u6557\uff1a"))
        finally:
            root.update()
            root.destroy()

    def test_form_uses_chinese_labels(self) -> None:
        root = create_main_window()
        try:
            form = root.form

            self.assertIn("建立教學文件架構", form.subtitle_label.cget("text"))
            self.assertEqual(form.title_label.cget("text"), "PPT 標題")
            self.assertEqual(form.sections_label.cget("text"), "章節")
            self.assertEqual(form.content_pages_label.cget("text"), "選取章節的內文頁")
            self.assertEqual(form.preview_label.cget("text"), "投影片預覽")
            self.assertEqual(form.generate_button.cget("text"), "產生 PPTX")
            self.assertEqual(form.add_section_button.cget("text"), "新增章節")
            self.assertEqual(form.add_content_button.cget("text"), "新增內文頁")
            self.assertEqual(form.template_label.cget("text"), "樣板 PPTX")
            self.assertEqual(form.choose_template_button.cget("text"), "選擇樣板")
        finally:
            root.update()
            root.destroy()

    def test_form_left_controls_do_not_overlap_preview(self) -> None:
        root = create_main_window()
        try:
            root.geometry("900x650")
            form = root.form
            form.set_document_title("這是測試文件")
            form.add_section("一")
            form.add_section("二")
            form.add_section("三")
            root.update()

            left_right_edge = form.left_panel.winfo_rootx() + form.left_panel.winfo_width()
            preview_left_edge = form.preview.winfo_rootx()

            self.assertGreaterEqual(form.left_panel.winfo_width(), 340)
            self.assertLessEqual(left_right_edge + 8, preview_left_edge)
            self.assertLessEqual(
                form.content_pages_label.winfo_rootx() + form.content_pages_label.winfo_width(),
                left_right_edge,
            )
        finally:
            root.update()
            root.destroy()

    def test_form_add_buttons_remain_visible_at_reported_window_size(self) -> None:
        root = create_main_window()
        try:
            root.geometry("900x650")
            form = root.form
            form.set_document_title("這是測試文件")
            form.add_section("一")
            form.add_section("二")
            form.add_section("三")
            root.update()

            left_x = form.left_panel.winfo_rootx()
            left_right_edge = left_x + form.left_panel.winfo_width()

            for button in (form.add_section_button, form.add_content_button):
                button_left = button.winfo_rootx()
                button_right = button_left + button.winfo_width()

                self.assertGreaterEqual(button.winfo_width(), 120)
                self.assertGreaterEqual(button.winfo_height(), 20)
                self.assertGreaterEqual(button_left, left_x)
                self.assertLessEqual(button_right, left_right_edge)
        finally:
            root.update()
            root.destroy()


if __name__ == "__main__":
    unittest.main()
