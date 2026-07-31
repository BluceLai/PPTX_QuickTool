from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pptx.util import Inches, Pt

from .page_plan import PagePlan, PlannedPage


SLIDE_WIDTH = 12192000
SLIDE_HEIGHT = 6858000
DEFAULT_TEMPLATE_PATH = Path(__file__).resolve().parent / "assets" / "default-training-document-template.pptx"
RETURN_TO_CONTENTS_TEXT = "回主目錄"
TABLE_OF_CONTENTS_RETURN_TEXT = "目錄"
HYPERLINK_BLUE = RGBColor(0x00, 0x70, 0xC0)

LAYOUT_NAMES_BY_PAGE_KIND = {
    "cover": ("Title with picture",),
    "table_of_contents": ("Contents", "Text"),
    "section_start": ("Contents", "Text"),
    "content": ("Text",),
}


@dataclass(frozen=True)
class TemplatePresentation:
    presentation: Presentation
    agenda_paragraph_styles: dict[int, object]


def export_page_plan_to_pptx(plan: PagePlan, output_path: Path, template_path: Path | None = None) -> Path:
    template = _new_presentation_from_template(template_path)
    presentation = template.presentation
    slide_by_page_id = {}
    slide_pages = []
    for page in plan.pages:
        slide = presentation.slides.add_slide(_layout_for_page(presentation, page))
        slide_by_page_id[page.page_id] = slide
        slide_pages.append((slide, page))

    table_of_contents_slide = slide_by_page_id["table-of-contents"]
    for slide, page in slide_pages:
        _populate_slide(
            slide,
            page,
            plan,
            slide_by_page_id,
            table_of_contents_slide,
            template.agenda_paragraph_styles,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def _new_presentation_from_template(template_path: Path | None) -> TemplatePresentation:
    path = Path(template_path) if template_path is not None else DEFAULT_TEMPLATE_PATH
    if not path.exists():
        raise FileNotFoundError(f"Template PPTX was not found: {path}")

    presentation = Presentation(path)
    agenda_paragraph_styles = _extract_agenda_paragraph_styles(presentation)
    _remove_all_slides(presentation)
    return TemplatePresentation(
        presentation=presentation,
        agenda_paragraph_styles=agenda_paragraph_styles,
    )


def _remove_all_slides(presentation) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _layout_for_page(presentation, page: PlannedPage):
    wanted_names = LAYOUT_NAMES_BY_PAGE_KIND[page.kind]
    for name in wanted_names:
        for layout in presentation.slide_layouts:
            if layout.name == name:
                return layout
    return presentation.slide_layouts[0]


def _populate_slide(
    slide,
    page: PlannedPage,
    plan: PagePlan,
    slide_by_page_id,
    table_of_contents_slide,
    agenda_paragraph_styles: dict[int, object],
) -> None:
    _set_title(slide, page.title)
    if page.kind == "table_of_contents":
        _populate_table_of_contents(slide, plan, slide_by_page_id, agenda_paragraph_styles)
    elif page.kind == "section_start":
        _populate_section_agenda(slide, page, plan, table_of_contents_slide, agenda_paragraph_styles)
    elif page.kind == "content":
        _clear_body_placeholder(slide)


def _set_title(slide, text: str) -> None:
    shape = _placeholder(slide, PP_PLACEHOLDER.TITLE)
    if shape is None:
        shape = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.75))
    _set_single_run_text(shape.text_frame, text)


def _clear_body_placeholder(slide) -> None:
    shape = _placeholder(slide, PP_PLACEHOLDER.BODY)
    if shape is not None:
        shape.text_frame.clear()


def _set_single_run_text(text_frame, text: str):
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    return run


def _add_link_text(slide, text: str, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.runs[0].font.color.rgb = HYPERLINK_BLUE
    paragraph.runs[0].font.underline = True
    return box


def _populate_table_of_contents(slide, plan: PagePlan, slide_by_page_id, agenda_paragraph_styles: dict[int, object]) -> None:
    body_shape = _placeholder(slide, PP_PLACEHOLDER.BODY)
    if body_shape is None:
        _add_table_of_contents_text_boxes(slide, plan, slide_by_page_id)
        return

    text_frame = body_shape.text_frame
    text_frame.clear()
    _append_agenda_line(text_frame, TABLE_OF_CONTENTS_RETURN_TEXT, paragraph_styles=agenda_paragraph_styles)
    for entry in plan.table_of_contents_entries:
        run = _append_agenda_line(
            text_frame,
            entry.title,
            is_link=True,
            paragraph_styles=agenda_paragraph_styles,
        )
        _link_run_to_slide(slide, run, slide_by_page_id[entry.target_page_id])
        section_page = _page_for_id(plan, entry.target_page_id)
        for content_title in _content_titles_for_section(plan, section_page):
            _append_agenda_line(text_frame, content_title, level=1, paragraph_styles=agenda_paragraph_styles)


def _populate_section_agenda(
    slide,
    page: PlannedPage,
    plan: PagePlan,
    table_of_contents_slide,
    agenda_paragraph_styles: dict[int, object],
) -> None:
    body_shape = _placeholder(slide, PP_PLACEHOLDER.BODY)
    if body_shape is None:
        _add_return_link(slide, table_of_contents_slide)
        return

    current_section_index = _section_index_from_page_id(page.page_id)
    text_frame = body_shape.text_frame
    text_frame.clear()
    run = _append_agenda_line(
        text_frame,
        TABLE_OF_CONTENTS_RETURN_TEXT,
        is_link=True,
        paragraph_styles=agenda_paragraph_styles,
    )
    _link_run_to_slide(slide, run, table_of_contents_slide)
    for entry_index, entry in enumerate(plan.table_of_contents_entries, start=1):
        section_page = _page_for_id(plan, entry.target_page_id)
        run = _append_agenda_line(
            text_frame,
            entry.title,
            is_current=entry_index == current_section_index,
            paragraph_styles=agenda_paragraph_styles,
        )
        if entry_index == current_section_index:
            run.font.bold = True
        for content_title in _content_titles_for_section(plan, section_page):
            _append_agenda_line(text_frame, content_title, level=1, paragraph_styles=agenda_paragraph_styles)


def _add_table_of_contents_text_boxes(slide, plan: PagePlan, slide_by_page_id) -> None:
    for entry_index, entry in enumerate(plan.table_of_contents_entries, start=1):
        link = _add_link_text(
            slide,
            entry.title,
            Inches(1.15),
            Inches(1.65 + (entry_index - 1) * 0.48),
            Inches(10.4),
            Inches(0.4),
        )
        link.click_action.target_slide = slide_by_page_id[entry.target_page_id]


def _append_agenda_line(
    text_frame,
    text: str,
    level: int = 0,
    is_link: bool = False,
    is_current: bool = False,
    paragraph_styles: dict[int, object] | None = None,
):
    paragraph = text_frame.paragraphs[0] if _is_empty_text_frame(text_frame) else text_frame.add_paragraph()
    paragraph.level = level
    _apply_paragraph_style(paragraph, level, paragraph_styles or {})
    run = paragraph.add_run()
    run.text = text
    if is_link:
        run.font.color.rgb = HYPERLINK_BLUE
        run.font.underline = True
    if is_current:
        run.font.bold = True
    return run


def _extract_agenda_paragraph_styles(presentation: Presentation) -> dict[int, object]:
    styles = _agenda_paragraph_styles_from_sample_slides(presentation)
    fallback_styles = _agenda_paragraph_styles_from_master(presentation)
    return {level: styles.get(level) or fallback_styles[level] for level in (0, 1) if level in styles or level in fallback_styles}


def _agenda_paragraph_styles_from_sample_slides(presentation: Presentation) -> dict[int, object]:
    styles = {}
    for slide in presentation.slides:
        if slide.slide_layout.name not in {"Contents", "Text"}:
            continue
        body_shape = _placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_shape is None:
            continue
        for paragraph in body_shape.text_frame.paragraphs:
            if not paragraph.text.strip() or paragraph._p.pPr is None:
                continue
            level = paragraph.level or 0
            styles.setdefault(level, deepcopy(paragraph._p.pPr))
        if 0 in styles and 1 in styles:
            return styles
    return styles


def _agenda_paragraph_styles_from_master(presentation: Presentation) -> dict[int, object]:
    styles = {}
    for master in presentation.slide_masters:
        for level in (0, 1):
            style = _master_body_level_style(master, level)
            if style is not None:
                styles.setdefault(level, style)
    return styles


def _master_body_level_style(master, level: int):
    tag_name = f".//p:txStyles/p:bodyStyle/a:lvl{level + 1}pPr"
    elements = master._element.xpath(tag_name)
    if not elements:
        return None
    style = deepcopy(elements[0])
    style.tag = qn("a:pPr")
    if level > 0:
        style.set("lvl", str(level))
    return style


def _apply_paragraph_style(paragraph, level: int, paragraph_styles: dict[int, object]) -> None:
    style = paragraph_styles.get(level)
    if style is None:
        return
    if paragraph._p.pPr is not None:
        paragraph._p.remove(paragraph._p.pPr)
    ppr = deepcopy(style)
    if level > 0:
        ppr.set("lvl", str(level))
    elif "lvl" in ppr.attrib:
        del ppr.attrib["lvl"]
    paragraph._p.insert(0, ppr)


def _is_empty_text_frame(text_frame) -> bool:
    return len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text


def _placeholder(slide, placeholder_type):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == placeholder_type:
            return shape
    return None


def _link_run_to_slide(slide, run, target_slide) -> None:
    relationship_id = slide.part.relate_to(target_slide.part, RT.SLIDE)
    run_properties = run._r.get_or_add_rPr()
    hyperlink = OxmlElement("a:hlinkClick")
    hyperlink.set(qn("r:id"), relationship_id)
    hyperlink.set("action", "ppaction://hlinksldjump")
    run_properties.append(hyperlink)


def _page_for_id(plan: PagePlan, page_id: str) -> PlannedPage:
    for page in plan.pages:
        if page.page_id == page_id:
            return page
    raise ValueError(f"Unknown page id: {page_id}")


def _content_titles_for_section(plan: PagePlan, section_page: PlannedPage) -> list[str]:
    section_index = _section_index_from_page_id(section_page.page_id)
    prefix = f"section-{section_index}-content-"
    return [page.title for page in plan.pages if page.page_id.startswith(prefix)]


def _section_index_from_page_id(page_id: str) -> int:
    return int(page_id.split("-")[1])


def _add_return_link(slide, table_of_contents_slide) -> None:
    link = _add_link_text(
        slide,
        RETURN_TO_CONTENTS_TEXT,
        Inches(10.0),
        Inches(6.55),
        Inches(2.3),
        Inches(0.35),
    )
    link.text_frame.paragraphs[0].font.size = Pt(12)
    link.click_action.target_slide = table_of_contents_slide
