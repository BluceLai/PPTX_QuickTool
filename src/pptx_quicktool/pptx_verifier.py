from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from .page_plan import PagePlan
from .pptx_exporter import (
    DEFAULT_TEMPLATE_PATH,
    RETURN_TO_CONTENTS_TEXT,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    TABLE_OF_CONTENTS_RETURN_TEXT,
)


@dataclass(frozen=True)
class VerificationResult:
    messages: list[str]

    @property
    def is_valid(self) -> bool:
        return len(self.messages) == 0


def verify_pptx_output(path: Path, plan: PagePlan, template_path: Path | None = None) -> VerificationResult:
    try:
        presentation = Presentation(path)
    except Exception:
        return VerificationResult(messages=["Output file cannot be opened as a PPTX."])

    messages: list[str] = []
    expected_width, expected_height = _expected_slide_size(template_path)
    if presentation.slide_width != expected_width or presentation.slide_height != expected_height:
        messages.append("Output slide size does not match the expected template size.")

    expected_slide_count = len(plan.pages)
    actual_slide_count = len(presentation.slides)
    if actual_slide_count != expected_slide_count:
        messages.append(f"Expected {expected_slide_count} slides, found {actual_slide_count}.")

    if actual_slide_count == expected_slide_count:
        messages.extend(_verify_required_text(presentation, plan))
        messages.extend(_verify_navigation_links(presentation, plan))

    return VerificationResult(messages=messages)


def _verify_required_text(presentation: Presentation, plan: PagePlan) -> list[str]:
    messages: list[str] = []
    for index, page in enumerate(plan.pages):
        text = _slide_text(presentation.slides[index])
        if page.title not in text:
            messages.append(f"Slide {index + 1} is missing required text: {page.title}.")
        if page.kind == "table_of_contents":
            for entry_index, entry in enumerate(plan.table_of_contents_entries, start=1):
                expected = f"{entry_index}. {entry.title}"
                if expected not in text:
                    messages.append(f"Slide {index + 1} is missing table of contents entry: {expected}.")
        if page.kind == "section_start" and TABLE_OF_CONTENTS_RETURN_TEXT not in text:
            messages.append(f"Slide {index + 1} is missing the table-of-contents return text.")
        if page.kind == "content" and RETURN_TO_CONTENTS_TEXT not in text:
            messages.append(f"Slide {index + 1} is missing the return-to-contents link text.")
    return messages


def _verify_navigation_links(presentation: Presentation, plan: PagePlan) -> list[str]:
    messages: list[str] = []
    slide_index_by_page_id = {page.page_id: index for index, page in enumerate(plan.pages)}
    table_of_contents_index = slide_index_by_page_id["table-of-contents"]
    table_of_contents_slide = presentation.slides[table_of_contents_index]
    table_of_contents_targets = _linked_slide_indices(presentation, table_of_contents_slide)
    for entry_index, entry in enumerate(plan.table_of_contents_entries, start=1):
        expected_index = slide_index_by_page_id[entry.target_page_id]
        if expected_index not in table_of_contents_targets:
            messages.append(f"Table of contents entry {entry_index} links to the wrong slide.")

    for slide_index, page in enumerate(plan.pages):
        if page.kind not in {"section_start", "content"}:
            continue
        linked_indices = _linked_slide_indices(presentation, presentation.slides[slide_index])
        if table_of_contents_index not in linked_indices:
            messages.append(f"Slide {slide_index + 1} return link points to the wrong slide.")

    return messages


def _expected_slide_size(template_path: Path | None) -> tuple[int, int]:
    path = template_path or DEFAULT_TEMPLATE_PATH
    if path.exists():
        presentation = Presentation(path)
        return presentation.slide_width, presentation.slide_height
    return SLIDE_WIDTH, SLIDE_HEIGHT


def _slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().replace("\v", "\n"))
    return "\n".join(texts)


def _linked_slide_indices(presentation: Presentation, slide) -> list[int]:
    indices = []
    for shape in slide.shapes:
        target_slide = shape.click_action.target_slide
        if target_slide is not None:
            indices.append(_target_slide_index(presentation, target_slide))
        for hyperlink in shape._element.xpath(".//a:rPr/a:hlinkClick"):
            relationship_id = hyperlink.get(qn("r:id"))
            if relationship_id:
                indices.append(_target_slide_part_index(presentation, slide.part.rels[relationship_id].target_part))
    return indices


def _target_slide_index(presentation: Presentation, target_slide) -> int:
    for index, slide in enumerate(presentation.slides):
        if slide == target_slide:
            return index
    raise ValueError("Target slide was not found in presentation.")


def _target_slide_part_index(presentation: Presentation, target_part) -> int:
    for index, slide in enumerate(presentation.slides):
        if slide.part == target_part:
            return index
    raise ValueError("Target slide part was not found in presentation.")
